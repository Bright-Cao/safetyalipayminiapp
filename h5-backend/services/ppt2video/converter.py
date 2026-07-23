import os
import sys
import argparse
import asyncio
import subprocess
import shutil
import tempfile
import json
from pptx import Presentation
from pdf2image import convert_from_path
import edge_tts
from moviepy import ImageClip, AudioFileClip, VideoFileClip, CompositeVideoClip, concatenate_videoclips
from PIL import Image

def report_progress(percent, status):
    """Outputs JSON progress line for Node.js process listener."""
    print(json.dumps({"progress": percent, "status": status}), flush=True)

def extract_notes(pptx_path):
    """Extracts speaker notes from PPTX slides."""
    prs = Presentation(pptx_path)
    notes = []
    for idx, slide in enumerate(prs.slides):
        note_text = ""
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            note_text = slide.notes_slide.notes_text_frame.text.strip()
        if not note_text:
            title = f"第{idx+1}页"
            if slide.shapes.title and slide.shapes.title.text:
                title = slide.shapes.title.text.strip()
            note_text = title
        notes.append(note_text)
    return notes

def extract_all_gifs(pptx_path, output_dir):
    """Scans PPTX slides and extracts animated GIF files with coordinates."""
    gif_data = []
    try:
        prs = Presentation(pptx_path)
        slide_w = prs.slide_width
        slide_h = prs.slide_height

        for idx, slide in enumerate(prs.slides):
            slide_gifs = []
            gif_count = 0
            for shape in slide.shapes:
                if shape.shape_type == 13: # PICTURE
                    try:
                        if hasattr(shape, 'image') and getattr(shape.image, 'content_type', '') == 'image/gif':
                            gif_filename = f"slide_{idx+1}_gif_{gif_count}.gif"
                            gif_path = os.path.join(output_dir, gif_filename)
                            with open(gif_path, 'wb') as f:
                                f.write(shape.image.blob)

                            # Verify animated
                            is_animated = False
                            try:
                                with Image.open(gif_path) as img:
                                    is_animated = hasattr(img, 'n_frames') and img.n_frames > 1
                            except Exception:
                                is_animated = True

                            if is_animated:
                                slide_gifs.append({
                                    'gif_path': gif_path,
                                    'left': shape.left / slide_w,
                                    'top': shape.top / slide_h,
                                    'width': shape.width / slide_w,
                                    'height': shape.height / slide_h
                                })
                                gif_count += 1
                            else:
                                os.remove(gif_path)
                    except Exception:
                        pass
            gif_data.append(slide_gifs)
    except Exception as e:
        print(f"GIF extract error: {e}")
    return gif_data

def convert_pptx_to_pdf(pptx_path, output_dir):
    """Converts PPTX to PDF using LibreOffice headless mode."""
    report_progress(15, "正在使用 LibreOffice 转换为 PDF...")
    soffice_bin = shutil.which("soffice") or shutil.which("libreoffice") or "/usr/bin/soffice"
    if not os.path.exists(soffice_bin) and not shutil.which(soffice_bin):
        raise RuntimeError("LibreOffice (soffice) 未在服务器上找到")

    cmd = [soffice_bin, "--headless", "--convert-to", "pdf", "--outdir", output_dir, pptx_path]
    res = subprocess.run(cmd, capture_output=True, text=True, timeout=180)
    if res.returncode != 0:
        raise RuntimeError(f"LibreOffice 转换失败: {res.stderr}")

    base_name = os.path.splitext(os.path.basename(pptx_path))[0]
    pdf_path = os.path.join(output_dir, base_name + ".pdf")
    if not os.path.exists(pdf_path):
        for f in os.listdir(output_dir):
            if f.endswith(".pdf"):
                return os.path.join(output_dir, f)
        raise FileNotFoundError("未查找到生成的 PDF 文件")
    return pdf_path

def convert_pdf_to_images(pdf_path, output_dir):
    """Converts PDF pages to PNG images using pdf2image (dpi=120)."""
    report_progress(35, "正在将 PDF 渲染为高清图片...")
    images = convert_from_path(pdf_path, dpi=120)
    image_paths = []
    for i, img in enumerate(images):
        img_path = os.path.join(output_dir, f"slide_{i+1}.png")
        img.save(img_path, "PNG")
        image_paths.append(img_path)
    return image_paths

async def generate_audio_files(notes, output_dir, voice="zh-CN-XiaoxiaoNeural", rate="+0%"):
    """Generates MP3 audio files for each slide note using edge-tts."""
    report_progress(55, "正在生成语音配音...")
    audio_paths = []
    total = len(notes)
    for i, note in enumerate(notes):
        audio_path = os.path.join(output_dir, f"audio_{i+1}.mp3")
        text = note if note.strip() else f"第 {i+1} 页幻灯片"
        communicate = edge_tts.Communicate(text, voice, rate=rate)
        await communicate.save(audio_path)
        audio_paths.append(audio_path)
        prog = 55 + int((i + 1) / total * 20)
        report_progress(prog, f"已生成第 {i+1}/{total} 页配音")
    return audio_paths

def compose_video(images, audio_files, gif_data, output_path):
    """Composes MP4 video from slide images, overlaying GIF animations if present."""
    report_progress(80, "正在合成 MP4 视频 (支持 GIF 动图叠加)...")
    clips = []
    for i, (img_path, audio_path) in enumerate(zip(images, audio_files)):
        audio_clip = AudioFileClip(audio_path)
        duration = max(audio_clip.duration, 1.0)
        bg_clip = ImageClip(img_path).with_duration(duration)
        w, h = bg_clip.w, bg_clip.h

        slide_gifs = gif_data[i] if i < len(gif_data) else []
        if slide_gifs:
            overlay_clips = [bg_clip]
            for gif in slide_gifs:
                try:
                    g_clip = VideoFileClip(gif['gif_path']).with_duration(duration)
                    if hasattr(g_clip, 'loop'):
                        g_clip = g_clip.loop(duration=duration)
                    x_pos = int(gif['left'] * w)
                    y_pos = int(gif['top'] * h)
                    gw = max(int(gif['width'] * w), 10)
                    gh = max(int(gif['height'] * h), 10)
                    g_clip = g_clip.resized((gw, gh)).with_position((x_pos, y_pos))
                    overlay_clips.append(g_clip)
                except Exception as e:
                    print(f"GIF overlay warning: {e}")
            img_clip = CompositeVideoClip(overlay_clips, size=(w, h)).with_duration(duration)
        else:
            img_clip = bg_clip

        img_clip = img_clip.with_audio(audio_clip)
        clips.append(img_clip)

    final_clip = concatenate_videoclips(clips, method="compose")
    final_clip.write_videofile(
        output_path,
        fps=24,
        codec="libx264",
        audio_codec="aac",
        preset="ultrafast",
        ffmpeg_params=["-threads", "1"],
        logger=None
    )

    final_clip.close()
    for c in clips:
        c.close()

    report_progress(100, "视频合成完成！")

async def process_ppt_to_video(pptx_path, output_path, voice="zh-CN-XiaoxiaoNeural", rate="+0%"):
    temp_dir = tempfile.mkdtemp(prefix="ppt2video_")
    try:
        report_progress(5, "正在解析 PPT 幻灯片与备注...")
        notes = extract_notes(pptx_path)
        if not notes:
            raise ValueError("PPT 文件中未找到幻灯片内容")

        report_progress(10, "正在检测与提取 PPT 中的 GIF 动态图片...")
        gif_data = extract_all_gifs(pptx_path, temp_dir)

        pdf_path = convert_pptx_to_pdf(pptx_path, temp_dir)
        images = convert_pdf_to_images(pdf_path, temp_dir)
        audio_files = await generate_audio_files(notes, temp_dir, voice=voice, rate=rate)

        if len(images) != len(audio_files):
            raise ValueError("幻灯片页数与音频数量不一致")

        compose_video(images, audio_files, gif_data, output_path)

    finally:
        shutil.rmtree(temp_dir, ignore_errors=True)

def main():
    parser = argparse.ArgumentParser(description="Linux PPT to Video Converter Service (with GIF support)")
    parser.add_argument("pptx_path", help="输入 PPTX 文件路径")
    parser.add_argument("--output", required=True, help="输出 MP4 视频路径")
    parser.add_argument("--voice", default="zh-CN-XiaoxiaoNeural", help="TTS 发音人")
    parser.add_argument("--rate", default="+0%", help="语速调整")

    args = parser.parse_args()

    try:
        asyncio.run(process_ppt_to_video(args.pptx_path, args.output, args.voice, args.rate))
    except Exception as e:
        print(json.dumps({"error": str(e)}), flush=True)
        sys.exit(1)

if __name__ == "__main__":
    main()
